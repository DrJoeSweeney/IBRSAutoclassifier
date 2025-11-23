"""
Document parsing module
Extracts text from various document formats
"""

import io
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from . import config


class DocumentParser:
    """Parse various document formats and extract text"""

    @staticmethod
    def detect_format(mime_type, filename=None):
        """
        Detect document format from MIME type

        Returns: format string ('pdf', 'docx', 'pptx', 'txt', 'image') or None
        """
        return config.SUPPORTED_MIME_TYPES.get(mime_type)

    @staticmethod
    def extract_text(file_content, mime_type, filename=None):
        """
        Extract text from document

        Args:
            file_content: Binary file content
            mime_type: MIME type of document
            filename: Original filename (optional)

        Returns:
            Extracted text as string

        Raises:
            ValueError: If format is unsupported or extraction fails
        """
        doc_format = DocumentParser.detect_format(mime_type)

        if not doc_format:
            raise ValueError(f"Unsupported MIME type: {mime_type}")

        try:
            if doc_format == 'pdf':
                return DocumentParser._extract_from_pdf(file_content)
            elif doc_format == 'docx':
                return DocumentParser._extract_from_docx(file_content)
            elif doc_format == 'pptx':
                return DocumentParser._extract_from_pptx(file_content)
            elif doc_format == 'txt':
                return DocumentParser._extract_from_text(file_content)
            elif doc_format == 'image':
                return DocumentParser._extract_from_image(file_content)
            else:
                raise ValueError(f"Unsupported format: {doc_format}")

        except Exception as e:
            raise ValueError(f"Failed to extract text from {doc_format}: {str(e)}")

    @staticmethod
    def _extract_from_pdf(file_content):
        """Extract text from PDF using PyPDF2 and pdfplumber as fallback"""
        text_parts = []

        try:
            # Try PyPDF2 first (faster)
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            # If we got reasonable text, return it
            combined_text = '\n'.join(text_parts)
            if len(combined_text.strip()) > 100:
                return combined_text

        except Exception as e:
            print(f"PyPDF2 extraction failed: {str(e)}, trying pdfplumber")

        # Fallback to pdfplumber for complex PDFs
        try:
            text_parts = []
            pdf_file = io.BytesIO(file_content)
            with pdfplumber.open(pdf_file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            return '\n'.join(text_parts)

        except Exception as e:
            raise ValueError(f"PDF extraction failed with both methods: {str(e)}")

    @staticmethod
    def _extract_from_docx(file_content):
        """Extract text from Word document"""
        doc_file = io.BytesIO(file_content)
        doc = Document(doc_file)

        text_parts = []

        # Extract from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)

        return '\n'.join(text_parts)

    @staticmethod
    def _extract_from_pptx(file_content):
        """Extract text from PowerPoint presentation"""
        ppt_file = io.BytesIO(file_content)
        prs = Presentation(ppt_file)

        text_parts = []

        for slide in prs.slides:
            # Extract from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)

            # Extract from notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text_parts.append(notes_text)

        return '\n'.join(text_parts)

    @staticmethod
    def _extract_from_text(file_content):
        """Extract text from plain text file"""
        try:
            # Try UTF-8 first
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1
            try:
                return file_content.decode('latin-1')
            except Exception as e:
                raise ValueError(f"Unable to decode text file: {str(e)}")

    @staticmethod
    def _extract_from_image(file_content):
        """Extract text from image using OCR"""
        try:
            # Open image
            image = Image.open(io.BytesIO(file_content))

            # Convert to grayscale for better OCR
            if image.mode != 'L':
                image = image.convert('L')

            # Apply OCR
            text = pytesseract.image_to_string(image)

            if not text.strip():
                raise ValueError("No text detected in image")

            return text

        except Exception as e:
            raise ValueError(f"OCR extraction failed: {str(e)}")

    @staticmethod
    def validate_extracted_text(text):
        """
        Validate that extracted text is sufficient for classification

        Args:
            text: Extracted text string

        Returns:
            tuple: (is_valid: bool, error_message: str or None)
        """
        if not text:
            return False, "No text extracted from document"

        text = text.strip()

        if len(text) < config.MIN_TEXT_LENGTH:
            return False, f"Extracted text too short ({len(text)} chars, minimum {config.MIN_TEXT_LENGTH})"

        # Check if text is mostly meaningful (not just whitespace/special chars)
        alphanumeric_count = sum(c.isalnum() for c in text)
        if alphanumeric_count < (len(text) * 0.3):
            return False, "Extracted text appears to contain mostly non-text characters"

        return True, None
